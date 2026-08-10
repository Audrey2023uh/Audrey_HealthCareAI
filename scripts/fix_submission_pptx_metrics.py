"""Build submission PPTX: reordered slides + accurate top-k metric wording."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

SRC_REORDERED = Path(
    r"C:\Users\audre\OneDrive\0-1- Urgent\2- Milstone 2\Milestone2_HealthcareAI_V2_REORDERED.pptx"
)
REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_REPO = REPO_ROOT / "Milestone2_HealthcareAI_V2.pptx"
OUT_ONEDRIVE = Path(
    r"C:\Users\audre\OneDrive\0-1- Urgent\2- Milstone 2\Milestone2_HealthcareAI_V2_SUBMISSION.pptx"
)
OUT_ONEDRIVE_REORDERED = Path(
    r"C:\Users\audre\OneDrive\0-1- Urgent\2- Milstone 2\Milestone2_HealthcareAI_V2_REORDERED.pptx"
)


def first_title(slide) -> str:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip().split("\n")[0][:120]
    return "(no text)"


def fix_metric_wording(prs: Presentation) -> list[str]:
    changes: list[str] = []
    for slide in prs.slides:
        if not first_title(slide).startswith("Evaluation"):
            continue
        for shape in slide.shapes:
            if not (hasattr(shape, "text") and shape.text and "Chunks / query" in shape.text):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == "Chunks / query":
                        run.text = "Configured top-k"
                        changes.append("label -> Configured top-k")
                    elif "Top-k after" in run.text:
                        run.text = "Configured default\n(eval returned 4–5)"
                        changes.append("subtitle -> Configured default (eval returned 4–5)")
            # value '5' stays — it is the configured TOP_K
            changes.append(f"final_text={shape.text!r}")
    return changes


def main() -> None:
    assert SRC_REORDERED.exists(), SRC_REORDERED
    shutil.copy2(SRC_REORDERED, OUT_REPO)

    prs = Presentation(str(OUT_REPO))
    changes = fix_metric_wording(prs)
    assert changes, "Expected to find Chunks / query metric card"
    prs.save(str(OUT_REPO))

    shutil.copy2(OUT_REPO, OUT_ONEDRIVE)
    shutil.copy2(OUT_REPO, OUT_ONEDRIVE_REORDERED)

    prs2 = Presentation(str(OUT_REPO))
    print(f"SLIDE_COUNT={len(prs2.slides)}")
    print("ORDER:")
    for i, s in enumerate(prs2.slides, 1):
        print(f"{i:02d}. {first_title(s)}")
    print("METRIC_CHANGES:")
    for c in changes:
        print(" ", c)

    eval_blob = ""
    for s in prs2.slides:
        if first_title(s).startswith("Evaluation"):
            eval_blob = "\n".join(sh.text for sh in s.shapes if hasattr(sh, "text") and sh.text)
            break
    print("EVAL_CHECK:")
    print(eval_blob)
    assert "1.00" in eval_blob
    assert "19 / 20" in eval_blob
    assert "Chunks / query" not in eval_blob
    assert "Configured top-k" in eval_blob
    assert "5" in eval_blob
    # Required early order
    titles = [first_title(s) for s in prs2.slides]
    assert titles[1].startswith("Resume Summary")
    assert "AI Tech Stack" in titles[2]
    assert titles[3].startswith("Data, Model")
    assert titles[4].startswith("Challenges")
    print(f"OUT_REPO={OUT_REPO}")
    print(f"OUT_SUBMISSION={OUT_ONEDRIVE}")
    print("OK")


if __name__ == "__main__":
    main()
