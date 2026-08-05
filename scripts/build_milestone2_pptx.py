"""
Build Milestone 2 PowerPoint — Healthcare AI CDSS (AI Engineering focus).

Reflects the ACTUAL implemented system in AI HEALTH CARE/:
  LangGraph 8-node CDSS graph, FAISS + TF-IDF RAG, Priority 1–7 ranking,
  Knowledge Update Agent, Streamlit UI, optional FastAPI / LLM.

Instructor requirement (AI Tech Stack.docx): yellow = implemented, green = ideal,
plus talk track for: what implemented / what ideal / why (trade-offs).

Inspired by Chip Huyen AI Engineering practice: implemented vs ideal honesty,
architecture layers, evaluation, deployment trade-offs.

Saves to:
  C:\\Users\\audre\\OneDrive\\1- Final project\\Milestone 2 - Healthcare AI CDSS Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(r"C:\Users\audre\OneDrive\1- Final project\Milestone 2 - Healthcare AI CDSS Presentation.pptx")

# ---- Palette (enterprise healthcare / Azure-architecture style) ----
NAVY = RGBColor(0x0F, 0x4C, 0x81)
DARK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFE, 0xF0, 0x8A)  # implemented
GREEN = RGBColor(0xBB, 0xF7, 0xD0)  # ideal
HEADER_BG = RGBColor(0x0F, 0x4C, 0x81)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
CYAN = RGBColor(0x0E, 0xA5, 0xE9)
TEAL = RGBColor(0x0D, 0x94, 0x88)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
SOFT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
SOFT_CYAN = RGBColor(0xCF, 0xFA, 0xFE)
SOFT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)
SOFT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
SOFT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
SOFT_ROSE = RGBColor(0xFF, 0xE4, 0xE6)
LAYER_BG = RGBColor(0xEE, 0xF2, 0xFF)


def set_run(run, text, size=14, bold=False, color=DARK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_footer(slide, prs, text="Milestone 2 · Audrey Rah · University of Houston · AI Engineering · Healthcare AI CDSS"):
    box = slide.shapes.add_textbox(Inches(0.4), prs.slide_height - Inches(0.38), Inches(12.5), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    set_run(p.add_run(), text, size=9, color=MUTED)


def title_bar(slide, prs, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.88))
    fill_shape(bar, HEADER_BG)
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.88), prs.slide_width, Inches(0.06))
    fill_shape(accent, CYAN)
    accent.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12.4), Inches(0.55))
    set_run(box.text_frame.paragraphs[0].add_run(), title, size=24, bold=True, color=WHITE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(1.05), Inches(12.4), Inches(0.32))
        set_run(sub.text_frame.paragraphs[0].add_run(), subtitle, size=12, color=MUTED)


def box(
    slide,
    left,
    top,
    width,
    height,
    text,
    fill=WHITE,
    border=BORDER,
    size=11,
    bold=True,
    color=DARK,
    align=PP_ALIGN.CENTER,
):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shp, fill)
    shp.line.color.rgb = border
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    # vertical-ish centering via empty top space is approximate
    run = tf.paragraphs[0].add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return shp


def arrow_right(slide, left, top, width=0.28, height=0.22, color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    fill_shape(shp, color)
    shp.line.fill.background()
    return shp


def arrow_down(slide, left, top, width=0.22, height=0.28, color=ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    fill_shape(shp, color)
    shp.line.fill.background()
    return shp


def set_cell_fill(cell, color: RGBColor):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def set_cell_text(cell, text, size=8, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)


def bullets_box(slide, left, top, width, height, items, size=12, fill=WHITE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(card, fill)
    card.line.color.rgb = BORDER
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), width - Inches(0.25), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        set_run(p.add_run(), f"• {t}", size=size, color=DARK)
    return card


def section_label(slide, left, top, width, text, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.32))
    fill_shape(shp, fill)
    shp.line.fill.background()
    tf = shp.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(tf.paragraphs[0].add_run(), text, size=11, bold=True, color=WHITE)
    return shp


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ==================================================================
    # 1 — Title
    # ==================================================================
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(bg, NAVY)
    bg.line.fill.background()
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.7), prs.slide_width, Inches(1.8))
    fill_shape(band, RGBColor(0x0B, 0x3A, 0x66))
    band.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.64), prs.slide_width, Inches(0.08))
    fill_shape(accent, CYAN)
    accent.line.fill.background()

    t = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(1.0))
    set_run(t.text_frame.paragraphs[0].add_run(), "Milestone 2", size=18, bold=True, color=CYAN)
    t2 = s.shapes.add_textbox(Inches(0.7), Inches(2.15), Inches(12), Inches(1.4))
    tf = t2.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0].add_run(),
        "Source-Linked Clinical Evidence CDSS",
        size=36,
        bold=True,
        color=WHITE,
    )
    t3 = s.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(12), Inches(1.2))
    tf = t3.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0].add_run(),
        "End-to-end AI Engineering Prototype  ·  RAG  ·  LangGraph Multi-Agent Workflow  ·  Evidence Hierarchy",
        size=16,
        color=SOFT_CYAN,
    )
    t4 = s.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(1.2))
    tf = t4.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "Audrey Rah  ·  University of Houston  ·  Houston, Texas, USA", size=15, bold=True, color=WHITE)
    p = tf.add_paragraph()
    set_run(p.add_run(), "Codebase: AI HEALTH CARE/  ·  Research/education prototype — not a medical device", size=12, color=SOFT_CYAN)

    # ==================================================================
    # 2 — Agenda
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(s, prs, "Agenda — AI Engineering Deliverable", "Aligned with proposal + AI Tech Stack.docx + Chip Huyen–style system thinking")
    items = [
        ("01", "What was built", "Resume-style system summary + honesty boundary"),
        ("02", "Architecture", "E2E, LangGraph graph, RAG, ranking, data flow"),
        ("03", "Agents & CDSS", "7-step clinical workflow + Knowledge Update"),
        ("04", "Tech stack", "Yellow = implemented · Green = ideal (required)"),
        ("05", "Eval & trade-offs", "Metrics, cost/complexity/scale, demo proof"),
    ]
    for i, (num, head, body) in enumerate(items):
        left = Inches(0.45 + (i % 5) * 2.55)
        top = Inches(1.55)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.4), Inches(4.6))
        fill_shape(card, WHITE)
        card.line.color.rgb = BORDER
        n = s.shapes.add_textbox(left + Inches(0.15), top + Inches(0.35), Inches(2.1), Inches(0.5))
        set_run(n.text_frame.paragraphs[0].add_run(), num, size=28, bold=True, color=CYAN)
        h = s.shapes.add_textbox(left + Inches(0.15), top + Inches(1.2), Inches(2.1), Inches(1.0))
        tf = h.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0].add_run(), head, size=16, bold=True, color=NAVY)
        b = s.shapes.add_textbox(left + Inches(0.15), top + Inches(2.4), Inches(2.1), Inches(1.8))
        tf = b.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0].add_run(), body, size=13, color=DARK)
    add_footer(s, prs)

    # ==================================================================
    # 3 — Resume bullets
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "What Was Built (Implemented — Not Proposed)",
        "Milestone 2 proof lives in AI HEALTH CARE/ · every claim maps to runnable code",
    )
    bullets = [
        (
            "End-to-end clinical decision-support prototype",
            "Streamlit clinician UI → LangGraph StateGraph (8 nodes) → source-linked recommendations with citations, confidence, and human-review flags.",
        ),
        (
            "Hierarchy-aware RAG over a public medical knowledge base",
            "FAISS IndexFlatIP + local TF-IDF embeddings; Priority 1–7 evidence ranking; superseded-guideline downgrade; optional OpenAI-compatible LLM with extractive offline fallback.",
        ),
        (
            "Explainable, measurable demo",
            "Agent trace, evidence tiles, conflict checks, citation precision eval script, optional FastAPI /ask, Knowledge Update Agent (URL probes + local supersession + index rebuild).",
        ),
    ]
    for i, (head, body) in enumerate(bullets):
        top = Inches(1.45 + i * 1.7)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), top, Inches(12.4), Inches(1.5))
        fill_shape(card, WHITE)
        card.line.color.rgb = BORDER
        accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), top, Inches(0.12), Inches(1.5))
        fill_shape(accent, ACCENT if i < 2 else TEAL)
        accent.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.85), top + Inches(0.25), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0].add_run(), f"{i+1}. {head}", size=17, bold=True, color=NAVY)
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        set_run(p.add_run(), body, size=13, color=DARK)
    add_footer(s, prs)

    # ==================================================================
    # 4 — Chip Huyen framing
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "AI Engineering Lens (Chip Huyen–style)",
        "Build for reliability under constraints — separate prototype reality from production ambition",
    )
    cols = [
        (
            "Application layer",
            SOFT_BLUE,
            [
                "Clinician workspace (Streamlit)",
                "Patient form + clinical query",
                "Citations + transparency panel",
                "Disclaimer + review flags",
            ],
        ),
        (
            "Orchestration layer",
            SOFT_CYAN,
            [
                "LangGraph StateGraph",
                "Shared GraphState",
                "Deterministic node edges",
                "Agent trace for audit",
            ],
        ),
        (
            "Model / RAG layer",
            SOFT_TEAL,
            [
                "Optional chat LLM",
                "Extractive fallback",
                "TF-IDF → FAISS retrieve",
                "Tier rerank + verify",
            ],
        ),
        (
            "Data / platform",
            SOFT_AMBER,
            [
                "Seed JSON medical KB",
                "Chunk store + FAISS",
                "Source registry probes",
                "Localhost runtime today",
            ],
        ),
    ]
    for i, (title, fill, items) in enumerate(cols):
        left = Inches(0.4 + i * 3.2)
        box(s, left, Inches(1.45), Inches(3.05), Inches(0.45), title, fill=NAVY, border=NAVY, size=13, color=WHITE)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.95), Inches(3.05), Inches(3.6))
        fill_shape(card, fill)
        card.line.color.rgb = BORDER
        tb = s.shapes.add_textbox(left + Inches(0.15), Inches(2.15), Inches(2.75), Inches(3.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, it in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            set_run(p.add_run(), f"• {it}", size=13, color=DARK)
    note = s.shapes.add_textbox(Inches(0.45), Inches(5.85), Inches(12.4), Inches(0.9))
    tf = note.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0].add_run(),
        "Engineering principle used in this talk: never present green-cell future components as if they already ship. "
        "Yellow cells are runnable today; green cells are the production-oriented roadmap.",
        size=13,
        bold=True,
        color=NAVY,
    )
    add_footer(s, prs)

    # ==================================================================
    # 5 — End-to-end architecture
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 1 — End-to-End System Architecture",
        "What a clinician request touches from UI to grounded answer (implemented path)",
    )
    # Layers as swimlanes
    layers = [
        (Inches(1.35), "Presentation", SOFT_BLUE, [("Streamlit CDSS UI", YELLOW), ("Optional FastAPI /ask", YELLOW)]),
        (Inches(2.55), "Orchestration", SOFT_CYAN, [("LangGraph CDSS StateGraph", YELLOW), ("GraphState + agent_trace", YELLOW)]),
        (
            Inches(3.75),
            "Intelligence",
            SOFT_TEAL,
            [("Query rewrite (opt LLM)", YELLOW), ("Verification + confidence", YELLOW), ("Recommendation builder", YELLOW)],
        ),
        (
            Inches(4.95),
            "Retrieval",
            SOFT_AMBER,
            [("TF-IDF embedder", YELLOW), ("FAISS IndexFlatIP", YELLOW), ("Priority 1–7 rerank", YELLOW)],
        ),
        (Inches(6.15), "Knowledge", SOFT_ROSE, [("seed_medical_kb.json", YELLOW), ("chunks.json + faiss.index", YELLOW), ("source_registry.json", YELLOW)]),
    ]
    for y, name, lane, comps in layers:
        lane_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), y, Inches(12.6), Inches(1.05))
        fill_shape(lane_bg, lane)
        lane_bg.line.color.rgb = BORDER
        lab = s.shapes.add_textbox(Inches(0.5), y + Inches(0.32), Inches(1.7), Inches(0.45))
        set_run(lab.text_frame.paragraphs[0].add_run(), name, size=12, bold=True, color=NAVY)
        for i, (txt, col) in enumerate(comps):
            box(
                s,
                Inches(2.4 + i * 3.4),
                y + Inches(0.2),
                Inches(3.15),
                Inches(0.65),
                txt,
                fill=col,
                border=NAVY,
                size=11,
            )
    # Flow label
    flow = s.shapes.add_textbox(Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.25))
    set_run(
        flow.text_frame.paragraphs[0].add_run(),
        "Request flow →  UI  →  LangGraph nodes  →  Retrieve+rank  →  Verify  →  Recommend  →  Transparency (citations)",
        size=11,
        bold=True,
        color=MUTED,
    )
    add_footer(s, prs)

    # ==================================================================
    # 6 — LangGraph execution graph
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 2 — LangGraph Execution Graph (Actual Edges)",
        "src/graph.py · StateGraph · entry = knowledge_update_agent · linear edges → END",
    )
    nodes = [
        ("Knowledge\nUpdate", SOFT_AMBER),
        ("Patient\nAssessment", SOFT_BLUE),
        ("Risk\nAnalysis", SOFT_CYAN),
        ("Query\nAgent", SOFT_TEAL),
        ("Evidence\nRetrieval", SOFT_PURPLE),
        ("Evidence\nVerification", SOFT_ROSE),
        ("Recommend +\nEvidence Summary", YELLOW),
        ("Transparency\n+ Final Answer", GREEN),
    ]
    # row 1: 0-3, row 2: 4-7
    for i, (label, fill) in enumerate(nodes):
        row = 0 if i < 4 else 1
        col = i if i < 4 else i - 4
        left = Inches(0.7 + col * 3.15)
        top = Inches(1.55 + row * 2.35)
        box(s, left, top, Inches(2.55), Inches(1.15), label, fill=fill, border=NAVY, size=13)
        if col < 3:
            arrow_right(s, left + Inches(2.6), top + Inches(0.45), width=0.4, height=0.25)
    # down arrow from row1 col3 to row2 col0 conceptually - show connector note
    arrow_down(s, Inches(11.5), Inches(2.8), width=0.28, height=0.55)
    note = s.shapes.add_textbox(Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.9))
    tf = note.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0].add_run(),
        "Compiled with g.compile(). Knowledge Update is optional per invoke (run_knowledge_update flag). "
        "All other nodes always execute. Shared GraphState carries patient, hits, verification, recommendation, timings, agent_trace.",
        size=12,
        color=DARK,
    )
    add_footer(s, prs)

    # ==================================================================
    # 7 — Multi-agent CDSS workflow
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 3 — Multi-Agent CDSS Workflow (Clinical Steps)",
        "Maps course CDSS requirements onto implemented LangGraph agents",
    )
    steps = [
        ("1", "Patient\nAssessment", "Form + text extract\nage/BP/DM/LDL/BMI"),
        ("2", "Risk\nAnalysis", "HTN stage, obesity,\nCV band, preventive"),
        ("3", "Evidence\nRetrieval", "FAISS top-k +\nPriority 1–7"),
        ("4", "Evidence\nVerification", "Conflicts, cites,\nconfidence"),
        ("5–6", "Recommend +\nSummary", "Lifestyle/Rx/screen\n+ source summary"),
        ("7", "Transparency", "Trace, sources,\nhuman-review flag"),
    ]
    for i, (n, title, body) in enumerate(steps):
        left = Inches(0.35 + i * 2.15)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.75), Inches(1.4), Inches(0.55), Inches(0.55))
        fill_shape(circ, NAVY)
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tf.paragraphs[0].add_run(), n, size=12, bold=True, color=WHITE)
        if i < len(steps) - 1:
            arrow_right(s, left + Inches(1.85), Inches(1.55), width=0.35, height=0.22)
        box(s, left, Inches(2.2), Inches(2.0), Inches(0.95), title, fill=SOFT_BLUE, border=NAVY, size=12)
        box(s, left, Inches(3.3), Inches(2.0), Inches(1.7), body, fill=WHITE, border=BORDER, size=11, bold=False)
    # persona callout
    box(
        s,
        Inches(0.4),
        Inches(5.35),
        Inches(12.5),
        Inches(1.35),
        "Persona continuity (Milestone 1): Dr. Lin — internal medicine resident needing source-linked evidence between patients.\n"
        "System supports judgment; licensed clinician makes the final decision. Prototype ≠ FDA medical device.",
        fill=LIGHT_BG,
        border=CYAN,
        size=13,
        bold=False,
        align=PP_ALIGN.LEFT,
    )
    add_footer(s, prs)

    # ==================================================================
    # 8 — RAG pipeline
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 4 — RAG Pipeline (Indexing + Query-Time)",
        "Classic RAG pattern with clinical ranking controls — not generation-only",
    )
    # Indexing path
    section_label(s, Inches(0.4), Inches(1.35), Inches(6.0), "Offline indexing path (scripts/build_index.py)")
    idx = [
        ("Seed KB\nJSON", SOFT_ROSE),
        ("Clean +\nChunk", SOFT_AMBER),
        ("TF-IDF\nfit/embed", SOFT_TEAL),
        ("FAISS\nIndexFlatIP", SOFT_CYAN),
        ("Persist\nindex+meta", YELLOW),
    ]
    for i, (t, c) in enumerate(idx):
        left = Inches(0.4 + i * 1.25)
        box(s, left, Inches(1.85), Inches(1.15), Inches(1.15), t, fill=c, border=NAVY, size=11)
        if i < len(idx) - 1:
            arrow_right(s, left + Inches(1.15), Inches(2.25), width=0.12, height=0.18, color=MUTED)

    # Query path
    section_label(s, Inches(6.9), Inches(1.35), Inches(6.0), "Online query path (retrieval_agent)")
    qry = [
        ("Clinical\nquery", SOFT_BLUE),
        ("Embed\nquery", SOFT_TEAL),
        ("ANN /\nIP search", SOFT_CYAN),
        ("Tier\nrerank", SOFT_PURPLE),
        ("Context\nblock", YELLOW),
    ]
    for i, (t, c) in enumerate(qry):
        left = Inches(6.9 + i * 1.2)
        box(s, left, Inches(1.85), Inches(1.1), Inches(1.15), t, fill=c, border=NAVY, size=11)
        if i < len(qry) - 1:
            arrow_right(s, left + Inches(1.1), Inches(2.25), width=0.12, height=0.18, color=MUTED)

    # Generation path
    section_label(s, Inches(0.4), Inches(3.4), Inches(12.5), "Grounded generation + verification")
    gen = [
        ("Retrieved\ncontext [n]", SOFT_CYAN),
        ("LLM chat\n(if key)", YELLOW),
        ("OR extractive\nfallback", YELLOW),
        ("Citation\nindex check", SOFT_ROSE),
        ("Confidence +\nreview flag", SOFT_AMBER),
        ("Final answer\n+ sources", GREEN),
    ]
    for i, (t, c) in enumerate(gen):
        left = Inches(0.45 + i * 2.1)
        box(s, left, Inches(3.95), Inches(1.95), Inches(1.2), t, fill=c, border=NAVY, size=12)
        if i < len(gen) - 1:
            arrow_right(s, left + Inches(1.95), Inches(4.4), width=0.18, height=0.22)

    note = s.shapes.add_textbox(Inches(0.45), Inches(5.5), Inches(12.4), Inches(1.2))
    tf = note.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0].add_run(),
        "Why RAG here: plain LLMs can invent guidelines. Retrieval-first + citation checks reduce unsupported claims. "
        "Current index: 19 docs / 20 chunks / embedder=local:tfidf (indexes/meta.json).",
        size=13,
        color=DARK,
    )
    add_footer(s, prs)

    # ==================================================================
    # 9 — Evidence hierarchy
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 5 — Evidence Hierarchy (Priority 1 → 7)",
        "src/evidence_rank.py · hard sort by clinical priority, then quality within tier — not similarity-only",
    )
    tiers = [
        ("P1", "Clinical Practice Guidelines", "USPSTF, ADA, ACC/AHA, NCCN, NICE, WHO…", RGBColor(0x1D, 0x4E, 0xD8), WHITE),
        ("P2", "Systematic Reviews / Meta-Analyses", "Cochrane-style / SR-MA summaries", RGBColor(0x25, 0x63, 0xEB), WHITE),
        ("P3", "Randomized Clinical Trials", "RCT abstracts / trial notes", RGBColor(0x3B, 0x82, 0xF6), WHITE),
        ("P4", "PubMed / PubMed Central", "Peer-reviewed literature style", RGBColor(0x60, 0xA5, 0xFA), DARK),
        ("P5", "AHRQ", "Gov. evidence / SDM themes", RGBColor(0x93, 0xC5, 0xFD), DARK),
        ("P6", "MedlinePlus", "Patient education (labeled)", RGBColor(0xBF, 0xDB, 0xFE), DARK),
        ("P7", "MedQuAD (testing only)", "Eval dataset — not guideline substitute", RGBColor(0xDB, 0xEA, 0xFE), DARK),
    ]
    for i, (p, name, detail, fill, tc) in enumerate(tiers):
        # trapezoid effect via shrinking width
        width = Inches(12.2 - i * 0.55)
        left = Inches(0.55 + i * 0.275)
        top = Inches(1.35 + i * 0.72)
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.65))
        fill_shape(shp, fill)
        shp.line.color.rgb = NAVY if i == 0 else BORDER
        tb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08), width - Inches(0.3), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0].add_run(), f"{p}  ·  {name}  —  {detail}", size=12, bold=(i < 3), color=tc)
    add_footer(s, prs)

    # ==================================================================
    # 10 — Retrieval & ranking pipeline
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 6 — Retrieval & Ranking Pipeline",
        "Hybrid: semantic candidates → clinical priority rerank → active-guideline preference",
    )
    pipeline = [
        ("1. Embed query", "TF-IDF vector\n(same space as index)", SOFT_TEAL),
        ("2. Candidate search", "FAISS IP top candidate_k\n(wide recall)", SOFT_CYAN),
        ("3. Annotate tiers", "infer Priority 1–7\n+ org / year / level", SOFT_BLUE),
        ("4. Rerank", "tier ↑ then quality\npenalize superseded", SOFT_PURPLE),
        ("5. Filter & cut", "score floor + top-k\nprefer non-superseded", SOFT_AMBER),
        ("6. Context pack", "format_evidence_block\n[1]…[k] for agents", YELLOW),
    ]
    for i, (h, b, c) in enumerate(pipeline):
        row = i // 3
        col = i % 3
        left = Inches(0.45 + col * 4.2)
        top = Inches(1.45 + row * 2.5)
        box(s, left, top, Inches(3.95), Inches(0.5), h, fill=NAVY, border=NAVY, size=13, color=WHITE)
        box(s, left, top + Inches(0.55), Inches(3.95), Inches(1.5), b, fill=c, border=BORDER, size=14, bold=False)
        if col < 2:
            arrow_right(s, left + Inches(3.95), top + Inches(0.95), width=0.25, height=0.25)
    add_footer(s, prs)

    # ==================================================================
    # 11 — Data flow / component interaction
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 7 — Component Interaction & Data Flow",
        "Modules in AI HEALTH CARE/src and how they call each other at runtime",
    )
    # Center orchestrator
    box(s, Inches(5.1), Inches(3.1), Inches(3.1), Inches(1.2), "LangGraph\nsrc/graph.py", fill=YELLOW, border=NAVY, size=14)
    # Surrounding
    comps = [
        (Inches(0.5), Inches(1.4), "app.py\nStreamlit UI", SOFT_BLUE),
        (Inches(4.0), Inches(1.4), "cdss.py\npatient/risk/recs", SOFT_CYAN),
        (Inches(7.5), Inches(1.4), "llm.py\nchat / extractive", SOFT_TEAL),
        (Inches(10.5), Inches(1.4), "api/main.py\nFastAPI /ask", SOFT_PURPLE),
        (Inches(0.5), Inches(5.1), "vectorstore.py\nFAISS search", SOFT_AMBER),
        (Inches(4.0), Inches(5.1), "evidence_rank.py\ntier rerank", SOFT_ROSE),
        (Inches(7.5), Inches(5.1), "embeddings.py\nTF-IDF", SOFT_TEAL),
        (Inches(10.5), Inches(5.1), "knowledge_update.py\nURL probes", SOFT_AMBER),
    ]
    for left, top, text, fill in comps:
        box(s, left, top, Inches(2.4), Inches(1.05), text, fill=fill, border=NAVY, size=12)
    # arrows from UI and to retrieval
    arrow_down(s, Inches(6.4), Inches(2.55), height=0.45)
    arrow_down(s, Inches(6.4), Inches(4.4), height=0.55)
    note = s.shapes.add_textbox(Inches(3.2), Inches(4.45), Inches(6.8), Inches(0.4))
    set_run(
        note.text_frame.paragraphs[0].add_run(),
        "invoke(question, patient_input, run_knowledge_update)  →  GraphState dict",
        size=11,
        bold=True,
        color=MUTED,
    )
    add_footer(s, prs)

    # ==================================================================
    # 12 — Knowledge update workflow
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 8 — Knowledge Update Workflow",
        "Implemented prototype loop — not full automated PDF ingestion of every society release",
    )
    ku = [
        ("Source registry\ndata/source_registry.json", SOFT_BLUE),
        ("HTTP probe\nstatus / Last-Modified / ETag", SOFT_CYAN),
        ("Compare to local\nseed doc versions", SOFT_TEAL),
        ("Apply supersession\nrules locally", SOFT_AMBER),
        ("Rebuild FAISS\n(optional / button)", YELLOW),
        ("Log report\nknowledge_update_log.json", GREEN),
    ]
    for i, (t, c) in enumerate(ku):
        left = Inches(0.4 + i * 2.15)
        box(s, left, Inches(1.7), Inches(2.0), Inches(2.0), t, fill=c, border=NAVY, size=12)
        if i < len(ku) - 1:
            arrow_right(s, left + Inches(2.0), Inches(2.5), width=0.18, height=0.25)
    bullets_box(
        s,
        Inches(0.4),
        Inches(4.2),
        Inches(6.1),
        Inches(2.4),
        [
            "Triggered from sidebar button or graph flag",
            "Probes USPSTF / ADA / ACC-AHA / MedlinePlus / AHRQ / PMC URLs",
            "Honest scope: URL freshness signals + local superseded_by",
        ],
        size=13,
    )
    bullets_box(
        s,
        Inches(6.8),
        Inches(4.2),
        Inches(6.1),
        Inches(2.4),
        [
            "Ideal (green / not claimed done): auto-ingest new guideline PDFs",
            "Ideal: continuous watchers + change detection + human approve",
            "Ideal: versioned corpus with audit trail for clinical governance",
        ],
        size=13,
        fill=SOFT_AMBER,
    )
    add_footer(s, prs)

    # ==================================================================
    # 13 — Deployment architecture
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Figure 9 — Deployment Architecture (Today vs Target)",
        "Yellow path = Milestone 2 demo runtime · Green path = production-oriented ideal",
    )
    # Today
    section_label(s, Inches(0.4), Inches(1.35), Inches(6.0), "Implemented deployment (local)")
    today = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.8), Inches(6.0), Inches(4.6))
    fill_shape(today, RGBColor(0xFF, 0xFB, 0xEB))
    today.line.color.rgb = ORANGE
    for i, t in enumerate(
        [
            "Developer laptop / classroom machine",
            "Python venv + requirements.txt",
            "streamlit run app.py  (UI)",
            "uvicorn api.main:app  (optional)",
            "indexes/ on local disk (FAISS + JSON)",
            "Optional .env OPENAI_API_KEY",
            "No PHI · public seed KB only",
        ]
    ):
        tb = s.shapes.add_textbox(Inches(0.65), Inches(2.05 + i * 0.55), Inches(5.5), Inches(0.45))
        set_run(tb.text_frame.paragraphs[0].add_run(), f"●  {t}", size=13, color=DARK)

    # Ideal
    section_label(s, Inches(6.9), Inches(1.35), Inches(6.0), "Ideal deployment (future)")
    ideal = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(6.0), Inches(4.6))
    fill_shape(ideal, RGBColor(0xEC, 0xFD, 0xF5))
    ideal.line.color.rgb = RGBColor(0x16, 0xA3, 0x4A)
    for i, t in enumerate(
        [
            "HIPAA-eligible cloud (Azure / AWS)",
            "Managed vector + hybrid search",
            "Azure OpenAI / Bedrock / Claude API",
            "React clinician console + auth",
            "CI eval + observability (LangSmith)",
            "Guideline watchers + approval gate",
            "EHR/FHIR optional with BAA controls",
        ]
    ):
        tb = s.shapes.add_textbox(Inches(7.15), Inches(2.05 + i * 0.55), Inches(5.5), Inches(0.45))
        set_run(tb.text_frame.paragraphs[0].add_run(), f"●  {t}", size=13, color=DARK)
    add_footer(s, prs)

    # ==================================================================
    # 14 — AI Tech Stack TABLE (instructor required)
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "AI Tech Stack — Implemented (Yellow) vs Ideal (Green)",
        "Required Milestone 2 format from AI Tech Stack.docx · answer the three instructor questions",
    )
    for i, (label, col) in enumerate(
        [("YELLOW = Implemented in Milestone 2", YELLOW), ("GREEN = Ideal / future architecture", GREEN)]
    ):
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45 + i * 4.3), Inches(1.25), Inches(4.0), Inches(0.3)
        )
        fill_shape(shp, col)
        shp.line.color.rgb = BORDER
        tb = s.shapes.add_textbox(Inches(0.55 + i * 4.3), Inches(1.27), Inches(3.8), Inches(0.26))
        set_run(tb.text_frame.paragraphs[0].add_run(), label, size=10, bold=True, color=DARK)

    rows, cols = 10, 7
    table_shape = s.shapes.add_table(rows, cols, Inches(0.3), Inches(1.65), Inches(12.7), Inches(4.55))
    table = table_shape.table
    headers = [
        "Layer (back → front)",
        "Really basic\n(local / simple)",
        "Freemium\n(hosted)",
        "Google (paid)",
        "AWS (paid)",
        "Microsoft (paid)",
        "Others (paid)",
    ]
    data = [
        [
            ("Data layer & vector store", "n"),
            ("JSON medical KB\n(local seed docs)", "y"),
            ("SQLite / Supabase\n(managed clinical KB)", "g"),
            ("BigQuery\nguideline warehouse", "g"),
            ("Amazon RDS +\nOpenSearch", "g"),
            ("Azure SQL +\nAI Search", "g"),
            ("FAISS local index\n(implemented)", "y"),
        ],
        [
            ("Model training & MLOps", "n"),
            ("VS Code + local\nPython scripts", "y"),
            ("Google Colab /\nW&B experiments", "g"),
            ("Vertex AI\npipelines", "g"),
            ("Amazon SageMaker", "g"),
            ("Azure Machine\nLearning", "g"),
            ("Databricks\neval harness", "g"),
        ],
        [
            ("Coding / IDE & agents", "n"),
            ("Terminal scripts", "y"),
            ("VS Code, Cursor", "y"),
            ("—", "n"),
            ("—", "n"),
            ("GitHub Copilot", "g"),
            ("LangGraph + LangChain\n(CrewAI / N8N ideal)", "y"),
        ],
        [
            ("Gen AI models", "n"),
            ("Ollama / local LLM\n(compatible path)", "y"),
            ("OpenAI-compatible API\n(OpenRouter optional)", "y"),
            ("Gemini API", "g"),
            ("Amazon Bedrock", "g"),
            ("Azure OpenAI", "g"),
            ("Anthropic Claude API\n(ideal clinical LLM)", "g"),
        ],
        [
            ("Frontend UI & dev tools", "n"),
            ("HTML/CSS prototypes", "n"),
            ("Streamlit\n(implemented UI)", "y"),
            ("Firebase Hosting", "g"),
            ("AWS Amplify", "g"),
            ("Azure Static\nWeb Apps", "g"),
            ("React / Next.js +\nFigma design system", "g"),
        ],
        [
            ("Hosting / runtime", "n"),
            ("Localhost\n(Streamlit / Uvicorn)", "y"),
            ("Render / Fly.io\nfree tier", "g"),
            ("Cloud Run", "g"),
            ("AWS Lambda +\nAPI Gateway", "g"),
            ("Azure App Service", "g"),
            ("HIPAA-eligible\nmanaged Kubernetes", "g"),
        ],
        [
            ("Embeddings &\nretrieval (bonus)", "n"),
            ("scikit-learn TF-IDF\n+ FAISS top-k", "y"),
            ("Managed vector\n(DB + hybrid search)", "g"),
            ("Vertex Matching\nEngine", "g"),
            ("OpenSearch k-NN\n+ re-ranker", "g"),
            ("Azure AI Search\nhybrid + semantic", "g"),
            ("Cohere / Voyage\nre-ranker (ideal)", "g"),
        ],
        [
            ("Orchestration &\nCDSS agents (bonus)", "n"),
            ("Python modules\n+ scripts", "y"),
            ("Hosted workflow\norchestrators", "g"),
            ("—", "n"),
            ("Step Functions\nagent workflows", "g"),
            ("Azure Logic Apps\n/ Durable Functions", "g"),
            ("LangGraph StateGraph\nCDSS graph (impl.)", "y"),
        ],
        [
            ("Eval / observability\n(bonus)", "n"),
            ("Local eval JSON\nlatency + citations", "y"),
            ("LangSmith free /\nopen telemetry", "g"),
            ("Cloud Monitoring", "g"),
            ("CloudWatch +\nX-Ray", "g"),
            ("Azure Monitor\nApp Insights", "g"),
            ("Continuous groundedness\n+ clinician override logs", "g"),
        ],
    ]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_fill(cell, HEADER_BG)
        set_cell_text(cell, h, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, row in enumerate(data):
        for j, (text, kind) in enumerate(row):
            cell = table.cell(i + 1, j)
            if kind == "y":
                set_cell_fill(cell, YELLOW)
            elif kind == "g":
                set_cell_fill(cell, GREEN)
            else:
                set_cell_fill(cell, LIGHT_BG)
            set_cell_text(
                cell,
                text,
                size=6 if j > 0 else 7,
                bold=(j == 0),
                color=DARK,
                align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER,
            )

    qa = s.shapes.add_textbox(Inches(0.35), Inches(6.35), Inches(12.6), Inches(0.7))
    tf = qa.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0].add_run(),
        "Instructor Qs: (1) Implemented stack? → yellow cells + Figures 1–2.  "
        "(2) Ideal stack? → green cells + Figure 9.  "
        "(3) Why? → next slide (cost, complexity, scalability, maintainability, latency).",
        size=11,
        bold=True,
        color=NAVY,
    )
    add_footer(s, prs)

    # ==================================================================
    # 15 — Trade-offs
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Why This Ideal Stack? — Engineering Trade-offs",
        "Answers AI Tech Stack.docx Q3 using Chip Huyen–style decision dimensions",
    )
    for j, h in enumerate(["Dimension", "Implemented choice (yellow)", "Ideal direction (green)"]):
        cell_shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35 + j * 4.25), Inches(1.3), Inches(4.15), Inches(0.4))
        fill_shape(cell_shp, NAVY)
        cell_shp.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.45 + j * 4.25), Inches(1.35), Inches(3.95), Inches(0.3))
        set_run(tb.text_frame.paragraphs[0].add_run(), h, size=12, bold=True, color=WHITE)

    rows_t = [
        ("Cost", "Local TF-IDF + optional LLM keeps demo near $0.", "Paid embeddings/LLM + cloud raise quality but add recurring cost / BAAs."),
        ("Complexity", "One Python repo + Streamlit is easy to run and grade.", "React + managed search + EHR FHIR expands integration surface."),
        ("Scalability", "FAISS-on-disk fine for 19-doc prototype KB.", "Managed vector + autoscaling APIs needed for multi-clinic load."),
        ("Maintainability", "Clear LangGraph nodes + seed JSON are inspectable.", "CI eval + guideline watchers reduce drift but need ops ownership."),
        ("Latency", "Local retrieval is fast; LLM optional.", "Re-rankers/LLMs add latency; cache + async jobs mitigate."),
        ("Safety / trust", "Citations, confidence, review flags, loud disclaimer.", "Clinical validation, audit logging, red-team eval harness."),
    ]
    for i, (dim, impl, ideal) in enumerate(rows_t):
        y = Inches(1.8 + i * 0.8)
        for j, (txt, bg) in enumerate([(dim, LIGHT_BG), (impl, YELLOW), (ideal, GREEN)]):
            shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35 + j * 4.25), y, Inches(4.15), Inches(0.72))
            fill_shape(shp, bg)
            shp.line.color.rgb = BORDER
            tb = s.shapes.add_textbox(Inches(0.45 + j * 4.25), y + Inches(0.1), Inches(3.95), Inches(0.55))
            tf = tb.text_frame
            tf.word_wrap = True
            set_run(tf.paragraphs[0].add_run(), txt, size=11, bold=(j == 0), color=DARK)
    add_footer(s, prs)

    # ==================================================================
    # 16 — Data, model, insights
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(s, prs, "Data, Model, and Key Insights", "Grounded in the running prototype — not aspirational slides")
    cols_info = [
        (
            "Data",
            [
                "Public-style seed KB (JSON): guideline summaries, MedlinePlus, AHRQ, PubMed/PMC-style, MedQuAD sample",
                "Pipeline: load → clean → chunk → TF-IDF → FAISS",
                "19 docs / 20 chunks indexed (meta.json)",
                "Knowledge Update URL registry + supersession",
                "Not claimed: bulk PMC crawl, EHR/FHIR PHI",
            ],
        ),
        (
            "Model / Agents",
            [
                "LangGraph StateGraph (8 nodes)",
                "Patient → Risk → Query → Retrieve → Verify → Recommend → Transparency",
                "Optional OpenAI-compatible chat LLM",
                "Offline extractive grounded fallback",
                "Streamlit UI + optional FastAPI",
            ],
        ),
        (
            "Key Insights",
            [
                "Guideline-first ranking beats pure semantic retrieval for clinical trust",
                "Superseded docs must be penalized or answers drift",
                "Citation precision = 1.0 on demo eval queries when answers cite [n]",
                "Offline path is required for reliable demos",
                "CDSS structure > plain chatbot Q&A for clinicians",
            ],
        ),
    ]
    for i, (title, items) in enumerate(cols_info):
        left = Inches(0.4 + i * 4.25)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(4.05), Inches(5.35))
        fill_shape(card, WHITE)
        card.line.color.rgb = BORDER
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), Inches(4.05), Inches(0.48))
        fill_shape(head, NAVY)
        head.line.fill.background()
        hb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.48), Inches(3.7), Inches(0.35))
        set_run(hb.text_frame.paragraphs[0].add_run(), title, size=16, bold=True, color=WHITE)
        body = s.shapes.add_textbox(left + Inches(0.18), Inches(2.05), Inches(3.7), Inches(4.5))
        tf = body.text_frame
        tf.word_wrap = True
        for j, it in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            set_run(p.add_run(), f"• {it}", size=12, color=DARK)
    add_footer(s, prs)

    # ==================================================================
    # 17 — Evaluation metrics
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Evaluation — Technical Proof Points",
        "scripts/evaluate.py → outputs/eval_results.json · proposal metrics mapped to what we can measure now",
    )
    metrics = [
        ("Citation precision", "1.00", "Valid [n] / total cites\n(demo eval set)"),
        ("Chunks / query", "5", "Top-k after\npriority rerank"),
        ("Indexed corpus", "19 / 20", "Documents / chunks\nin FAISS"),
        ("Retrieval path", "FAISS", "TF-IDF · hierarchy\nrerank"),
    ]
    for i, (label, value, hint) in enumerate(metrics):
        left = Inches(0.45 + i * 3.2)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(3.0), Inches(1.7))
        fill_shape(card, WHITE)
        card.line.color.rgb = BORDER
        tb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.5), Inches(2.7), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0].add_run(), label, size=12, bold=True, color=MUTED)
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), value, size=28, bold=True, color=NAVY)
        p3 = tf.add_paragraph()
        set_run(p3.add_run(), hint, size=11, color=MUTED)

    # Comparison table
    box(s, Inches(0.45), Inches(3.4), Inches(12.4), Inches(0.45), "RAG vs plain LLM (qualitative — matches proposal evaluation plan)", fill=NAVY, border=NAVY, size=14, color=WHITE)
    compare = [
        ("Dimension", "With RAG (implemented)", "Plain LLM (no retrieval)"),
        ("Source links", "Yes — chunk IDs, orgs, URLs", "None / invented references risk"),
        ("Groundedness", "Constrained to retrieved context", "Fluent but may hallucinate guidelines"),
        ("Audit trail", "Agent trace + timings", "No retrieval audit"),
        ("Offline demo", "Works via extractive fallback", "Requires API key"),
    ]
    for i, row in enumerate(compare):
        y = Inches(3.95 + i * 0.48)
        for j, txt in enumerate(row):
            bg = NAVY if i == 0 else (YELLOW if j == 1 and i > 0 else (SOFT_ROSE if j == 2 and i > 0 else LIGHT_BG))
            tc = WHITE if i == 0 else DARK
            box(
                s,
                Inches(0.45 + j * 4.15),
                y,
                Inches(4.05),
                Inches(0.45),
                txt,
                fill=bg,
                border=BORDER,
                size=11,
                bold=(i == 0 or j == 0),
                color=tc,
            )
    add_footer(s, prs)

    # ==================================================================
    # 18 — Challenges & learnings
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(s, prs, "Challenges and Learnings", "Engineering reality from building the prototype")
    left_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.35), Inches(6.1), Inches(5.4))
    fill_shape(left_card, WHITE)
    left_card.line.color.rgb = BORDER
    right_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.4))
    fill_shape(right_card, WHITE)
    right_card.line.color.rgb = BORDER
    lh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.35), Inches(6.1), Inches(0.48))
    fill_shape(lh, RGBColor(0xDC, 0x26, 0x26))
    lh.line.fill.background()
    rh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.35), Inches(6.1), Inches(0.48))
    fill_shape(rh, RGBColor(0x16, 0xA3, 0x4A))
    rh.line.fill.background()
    set_run(s.shapes.add_textbox(Inches(0.55), Inches(1.42), Inches(5.8), Inches(0.35)).text_frame.paragraphs[0].add_run(), "Challenges", size=16, bold=True, color=WHITE)
    set_run(s.shapes.add_textbox(Inches(6.95), Inches(1.42), Inches(5.8), Inches(0.35)).text_frame.paragraphs[0].add_run(), "Learnings", size=16, bold=True, color=WHITE)

    challenges = [
        "Heavy embedding stacks broke installs → defaulted to offline TF-IDF.",
        "Pure semantic search overweighted patient-education pages.",
        "Superseded guidelines retrieve unless ranking penalizes them.",
        "UI had to become clinician-readable without rewriting agents.",
        "Must label indexed/cached vs true live retrieval honestly.",
        "Medical prototypes need loud disclaimers + review flags.",
    ]
    learnings = [
        "Separate implemented vs ideal stack explicitly (AI engineering honesty).",
        "Evidence hierarchy + citation checks beat fluency-only answers.",
        "LangGraph boundaries make CDSS steps demoable and debuggable.",
        "Local-first design improves reliability and lowers cost.",
        "Trade-offs ARE the story: cost, latency, scale, maintainability.",
        "Clinician support ≠ replacement — transparency builds trust.",
    ]
    lb = s.shapes.add_textbox(Inches(0.6), Inches(2.05), Inches(5.7), Inches(4.5))
    tf = lb.text_frame
    tf.word_wrap = True
    for i, t in enumerate(challenges):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        set_run(p.add_run(), f"• {t}", size=13, color=DARK)
    rb = s.shapes.add_textbox(Inches(7.0), Inches(2.05), Inches(5.7), Inches(4.5))
    tf = rb.text_frame
    tf.word_wrap = True
    for i, t in enumerate(learnings):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        set_run(p.add_run(), f"• {t}", size=13, color=DARK)
    add_footer(s, prs)

    # ==================================================================
    # 19 — Demo checklist
    # ==================================================================
    s = prs.slides.add_slide(blank)
    title_bar(
        s,
        prs,
        "Live Demo — Proof of Implementation",
        "Camera ON · show running system + codebase path AI HEALTH CARE/",
    )
    demo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.35), Inches(8.3), Inches(5.4))
    fill_shape(demo, WHITE)
    demo.line.color.rgb = BORDER
    dh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.35), Inches(8.3), Inches(0.48))
    fill_shape(dh, NAVY)
    dh.line.fill.background()
    set_run(
        s.shapes.add_textbox(Inches(0.55), Inches(1.42), Inches(8.0), Inches(0.35)).text_frame.paragraphs[0].add_run(),
        "Demo script (record with camera ON)",
        size=14,
        bold=True,
        color=WHITE,
    )
    steps_demo = [
        "Start: python scripts\\build_index.py then streamlit run app.py (or START_DEMO.bat).",
        "Enter patient fields + clinical question (e.g., ASCVD statin intensity / A1C target).",
        "Run CDSS — show KPIs, assessment, risk, evidence tiles (guidelines first).",
        "Open verification conflicts, recommendations, citations, agent trace, confidence.",
        "Show Knowledge Update / Rebuild index buttons (maintenance).",
        "Point to src/graph.py, evidence_rank.py, vectorstore.py in the repo.",
        "State clearly: research prototype, not a medical device; clinician decides.",
    ]
    db = s.shapes.add_textbox(Inches(0.6), Inches(2.05), Inches(7.9), Inches(4.5))
    tf = db.text_frame
    tf.word_wrap = True
    for i, t in enumerate(steps_demo):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        set_run(p.add_run(), f"{i+1}. {t}", size=13, color=DARK)

    side = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(1.35), Inches(3.95), Inches(5.4))
    fill_shape(side, LIGHT_BG)
    side.line.color.rgb = BORDER
    stitle = s.shapes.add_textbox(Inches(9.1), Inches(1.55), Inches(3.6), Inches(0.4))
    set_run(stitle.text_frame.paragraphs[0].add_run(), "Honesty check", size=15, bold=True, color=NAVY)
    sbody = s.shapes.add_textbox(Inches(9.1), Inches(2.15), Inches(3.6), Inches(4.3))
    tf = sbody.text_frame
    tf.word_wrap = True
    for i, t in enumerate(
        [
            "Implemented: Streamlit, LangGraph CDSS graph, FAISS, TF-IDF, priority ranking, citations, FastAPI optional, Knowledge Update probes.",
            "Ideal (not done): managed vector DB, neural embeddings + re-ranker, HIPAA cloud, React app, full PMC/EHR.",
            "Folder: ...\\1- Final project\\AI HEALTH CARE",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        set_run(p.add_run(), f"• {t}", size=12, color=DARK)
    add_footer(s, prs)

    # ==================================================================
    # 20 — Closing
    # ==================================================================
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill_shape(bg, NAVY)
    bg.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0), prs.slide_width, Inches(0.1))
    fill_shape(accent, CYAN)
    accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2))
    set_run(t.text_frame.paragraphs[0].add_run(), "Closing", size=18, bold=True, color=CYAN)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.0))
    tf = t2.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0].add_run(),
        "Better clinical decisions need evidence clinicians can verify.\n"
        "This Milestone 2 system demonstrates AI engineering for that goal:\n"
        "RAG + agents + ranking + transparency — implemented, measured, and demoable.",
        size=20,
        bold=True,
        color=WHITE,
    )
    t3 = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5))
    tf = t3.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "Audrey Rah  ·  University of Houston  ·  Houston, Texas, USA", size=16, color=SOFT_CYAN)
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(), "Questions?", size=22, bold=True, color=WHITE)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    # also refresh the Copy if present
    copy = OUT.with_name(OUT.stem + " - Copy.pptx")
    try:
        prs.save(str(copy))
    except Exception:
        pass
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
