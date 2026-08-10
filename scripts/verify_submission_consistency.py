"""Consistency verification after Milestone 2 submission fixes."""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))


def main() -> None:
    # --- Graph: 8 nodes ---
    from src.graph import build_graph, run_pipeline
    import inspect
    from src import graph as gmod

    src = inspect.getsource(gmod.build_graph)
    nodes = re.findall(r'add_node\("([^"]+)"', src)
    assert len(nodes) == 8, nodes
    print("PASS 8-node LangGraph:", nodes)

    # --- Config / metrics files ---
    from src.config import TOP_K

    assert TOP_K == 5
    meta = json.loads((ROOT / "indexes" / "meta.json").read_text(encoding="utf-8"))
    assert meta["n_docs"] == 19 and meta["n_chunks"] == 20
    ev = json.loads((ROOT / "outputs" / "eval_results.json").read_text(encoding="utf-8"))
    chunks = [r["n_chunks"] for r in ev]
    precs = [r["citation_accuracy"]["citation_precision"] for r in ev]
    assert chunks == [4, 4, 5], chunks
    assert all(p == 1.0 for p in precs), precs
    print("PASS metrics files:", {"TOP_K": TOP_K, "meta": (meta["n_docs"], meta["n_chunks"]), "eval_chunks": chunks})

    # --- RAG smoke ---
    from src.vectorstore import get_store

    store = get_store()
    hits = store.search(
        "What A1C target is commonly used for nonpregnant adults with type 2 diabetes?",
        top_k=TOP_K,
    )
    assert hits and all("evidence_tier" in h for h in hits)
    out = run_pipeline(
        "What A1C target is commonly used for nonpregnant adults with type 2 diabetes?"
    )
    assert out.get("agent_trace")
    assert out.get("hits") is not None
    print("PASS RAG+pipeline:", {"n_hits": len(out.get("hits") or []), "trace": [t["agent"] for t in out["agent_trace"]]})

    # --- deps ---
    import fastapi
    import uvicorn

    print("PASS fastapi/uvicorn:", fastapi.__version__, uvicorn.__version__)

    # --- README ---
    readme = (ROOT / "README.md").read_text(encoding="utf-8")
    assert "4-agent" not in readme
    assert "8-node" in readme or "8 nodes" in readme
    assert "OneDrive\\1- Final project" not in readme
    assert "OPENROUTER_API_KEY" in readme
    assert "git clone" in readme
    print("PASS README wording")

    # --- requirements ---
    req = (ROOT / "requirements.txt").read_text(encoding="utf-8")
    assert "fastapi" in req and "uvicorn" in req
    print("PASS requirements.txt includes fastapi/uvicorn")

    # --- PPTX ---
    pptx_path = ROOT / "Milestone2_HealthcareAI_V2.pptx"
    prs = Presentation(str(pptx_path))
    titles = []
    for s in prs.slides:
        title = "(no text)"
        for sh in s.shapes:
            if hasattr(sh, "text") and sh.text.strip():
                title = sh.text.strip().split("\n")[0]
                break
        titles.append(title)
    assert len(titles) == 21
    assert titles[1].startswith("Resume Summary")
    assert "AI Tech Stack" in titles[2]
    assert titles[3].startswith("Data, Model")
    assert titles[4].startswith("Challenges")
    eval_blob = ""
    for s, t in zip(prs.slides, titles):
        if t.startswith("Evaluation"):
            eval_blob = "\n".join(sh.text for sh in s.shapes if hasattr(sh, "text") and sh.text)
    assert "Configured top-k" in eval_blob
    assert "Chunks / query" not in eval_blob
    assert "1.00" in eval_blob and "19 / 20" in eval_blob
    print("PASS PPTX order+metrics")
    print("ORDER:")
    for i, t in enumerate(titles, 1):
        print(f"  {i:02d}. {t}")

    # --- secrets scan (simple) ---
    bad = []
    for p in ROOT.rglob("*"):
        if ".git" in p.parts or p.suffix.lower() in {".png", ".index", ".pkl", ".pptx", ".docx"}:
            continue
        if p.is_file() and p.stat().st_size < 2_000_000:
            try:
                txt = p.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                continue
            if re.search(r"sk-(?:or-)?[a-zA-Z0-9]{20,}", txt) and "replace-me" not in txt and "your-openrouter" not in txt:
                # allow placeholder lists
                if "sk-or-v1-replace-me" in txt or "sk-xxx" in txt:
                    continue
                if re.search(r"sk-(?:or-v1-)?[a-f0-9]{32,}", txt):
                    bad.append(str(p))
    assert not bad, bad
    print("PASS no real secrets detected")
    print("ALL_CHECKS_PASSED")


if __name__ == "__main__":
    main()
